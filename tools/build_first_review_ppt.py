#!/usr/bin/env python3
"""Build the editable First Review deck and its rendered review copy.

The PowerPoint uses native text boxes and shapes so that the author can edit it.
The PDF and PNG contact sheet are generated from the same scene description for
review on machines without an Office renderer.
"""

from __future__ import annotations

from dataclasses import dataclass, field
from datetime import date
from pathlib import Path
import math
import shutil
import textwrap

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "deliverables" / "review-1"
PREVIEWS = OUT / "previews"

PPTX_PATH = OUT / "semantic-communication-first-review.pptx"
PDF_PATH = OUT / "semantic-communication-first-review.pdf"
CONTACT_PATH = OUT / "semantic-communication-first-review-contact-sheet.png"

SW, SH = 13.333, 7.5
PX_W, PX_H = 1600, 900

IVORY = "FFFFFF"
PAPER = "FFFFFF"
INK = "111111"
MUTED = "555555"
FAINT = "999999"
LINE = "D9D9D9"
NAVY = "111111"
BURGUNDY = "111111"
GREEN = "333333"
AMBER = "555555"
PALE_NAVY = "F3F3F3"
PALE_RED = "F3F3F3"
PALE_GREEN = "F3F3F3"
PALE_AMBER = "F3F3F3"
WHITE = "FFFFFF"

TITLE_FONT = "Arial"
BODY_FONT = "Arial"
MATH_FONT = "Arial"
MONO_FONT = "Courier New"

FONT_FILES = {
    # The review renderer is intentionally dependency-light. PowerPoint uses
    # the Office-facing font names above; local PNG/PDF proofs use the fonts
    # available on this WSL host.
    "serif": "/usr/share/fonts/Adwaita/AdwaitaSans-Regular.ttf",
    "serif_bold": "/usr/share/fonts/Adwaita/AdwaitaSans-Regular.ttf",
    "sans": "/usr/share/fonts/Adwaita/AdwaitaSans-Regular.ttf",
    "sans_bold": "/usr/share/fonts/Adwaita/AdwaitaSans-Regular.ttf",
    "mono": "/usr/share/fonts/Adwaita/AdwaitaMono-Regular.ttf",
}


def rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value)


def px(v: float, axis: str) -> int:
    return round(v * (PX_W / SW if axis == "x" else PX_H / SH))


def color(value: str) -> tuple[int, int, int]:
    return tuple(int(value[i : i + 2], 16) for i in (0, 2, 4))


@dataclass
class Element:
    kind: str
    x: float
    y: float
    w: float = 0
    h: float = 0
    text: str = ""
    size: float = 16
    fill: str | None = None
    stroke: str | None = None
    stroke_width: float = 1
    text_color: str = INK
    font: str = "sans"
    bold: bool = False
    italic: bool = False
    align: str = "left"
    valign: str = "top"
    radius: float = 0.08
    opacity: int = 255
    margin: float = 0.05
    rotation: float = 0


@dataclass
class SlideScene:
    number: int
    section: str
    title: str
    criteria: tuple[str, ...]
    citation: str = ""
    elements: list[Element] = field(default_factory=list)

    def text(self, x, y, w, h, text, **kwargs):
        self.elements.append(Element("text", x, y, w, h, text=text, **kwargs))

    def rect(self, x, y, w, h, **kwargs):
        self.elements.append(Element("rect", x, y, w, h, **kwargs))

    def line(self, x, y, w, h, **kwargs):
        self.elements.append(Element("line", x, y, w, h, **kwargs))

    def circle(self, x, y, w, h, **kwargs):
        self.elements.append(Element("circle", x, y, w, h, **kwargs))

    def arrow(self, x, y, w, h, **kwargs):
        self.elements.append(Element("arrow", x, y, w, h, **kwargs))


def add_header(scene: SlideScene) -> None:
    scene.text(0.58, 0.25, 4.6, 0.24, "CAPSTONE · FIRST REVIEW · 18–22 AUG 2026",
               size=8.5, text_color=MUTED, bold=True, font="sans")
    scene.text(11.65, 0.23, 1.05, 0.25, f"§ {scene.number:02d}", size=10,
               text_color=BURGUNDY, bold=True, align="right", font="mono")
    scene.text(0.58, 0.62, 11.9, 0.52, scene.title, size=27, bold=True,
               font="serif", text_color=INK, valign="mid")
    scene.line(0.58, 1.20, 12.15, 0, stroke=BURGUNDY, stroke_width=1.3)


def add_footer(scene: SlideScene) -> None:
    scene.line(0.58, 7.05, 12.15, 0, stroke=LINE, stroke_width=0.7)
    x = 0.58
    for criterion in scene.criteria:
        w = 0.28 + 0.064 * len(criterion)
        scene.rect(x, 7.13, w, 0.22, fill=PALE_NAVY, stroke=None, radius=0.05)
        scene.text(x + 0.07, 7.16, w - 0.14, 0.13, criterion.upper(), size=6.8,
                   text_color=NAVY, bold=True, font="sans", valign="mid")
        x += w + 0.10
    if scene.citation:
        scene.text(4.25, 7.13, 7.75, 0.20, scene.citation, size=6.4,
                   text_color=MUTED, font="sans", align="right", valign="mid")
    scene.text(12.18, 7.12, 0.52, 0.20, f"{scene.number:02d}/12", size=7.2,
               text_color=MUTED, font="mono", align="right", valign="mid")


def bullet(scene: SlideScene, x: float, y: float, w: float, text_value: str,
           *, size: float = 14, accent: str = BURGUNDY, lines: float = 0.52,
           text_color: str = INK) -> None:
    scene.rect(x, y + 0.13, 0.07, 0.07, fill=accent, stroke=None, radius=0)
    scene.text(x + 0.18, y, w - 0.18, lines, text_value, size=size,
               text_color=text_color, font="sans", valign="top")


def label(scene: SlideScene, x, y, w, text_value, *, fill=PALE_NAVY, ink=NAVY):
    scene.rect(x, y, w, 0.27, fill=fill, stroke=None, radius=0.04)
    scene.text(x + 0.08, y + 0.04, w - 0.16, 0.16, text_value.upper(), size=7.3,
               text_color=ink, bold=True, font="sans", valign="mid")


def card(scene: SlideScene, x, y, w, h, heading, body, *, accent=NAVY,
         body_size=12.5, fill=PAPER, tag=None):
    scene.rect(x, y, w, h, fill=fill, stroke=LINE, stroke_width=0.8, radius=0.08)
    scene.rect(x, y, 0.06, h, fill=accent, stroke=None, radius=0)
    if tag:
        label(scene, x + 0.24, y + 0.20, min(1.12, 0.35 + 0.075 * len(tag)), tag,
              fill=PALE_NAVY if accent == NAVY else PALE_RED,
              ink=accent)
        head_y = y + 0.60
    else:
        head_y = y + 0.23
    scene.text(x + 0.24, head_y, w - 0.46, 0.35, heading, size=15.5,
               bold=True, font="serif", text_color=INK)
    scene.text(x + 0.24, head_y + 0.48, w - 0.46, h - (head_y - y) - 0.65,
               body, size=body_size, font="sans", text_color=MUTED)


def build_scenes() -> list[SlideScene]:
    slides: list[SlideScene] = []

    # 1 — accessible title / thesis
    s = SlideScene(1, "Thesis", "Semantic Communication for Edge AI",
                   ("Motivation", "Objectives"),
                   "Official scope: task-oriented DJSCC for image classification · First Review baseline 12 Aug 2026")
    add_header(s)
    s.text(0.72, 1.48, 11.86, 0.45,
           "Can a machine communicate what matters for a task—without sending every pixel perfectly?",
           size=18.2, text_color=NAVY, font="serif", italic=True, align="center")
    # Simple application story: camera → constrained link → receiver AI.
    story = [
        (1.12, "EDGE CAMERA", "observes an image", NAVY, PALE_NAVY),
        (5.11, "LIMITED + NOISY LINK", "cannot carry everything reliably", BURGUNDY, PALE_RED),
        (9.10, "RECEIVER AI", "only needs the correct class", GREEN, PALE_GREEN),
    ]
    for x, heading, body, accent, pale in story:
        s.rect(x, 2.42, 3.10, 1.36, fill=pale, stroke=None, radius=0.08)
        s.text(x + 0.22, 2.76, 2.66, 0.23, heading, size=10.5, bold=True,
               font="sans", text_color=accent, align="center")
        s.text(x + 0.22, 3.23, 2.66, 0.24, body, size=11.2,
               font="serif", text_color=INK, align="center")
    s.arrow(4.32, 3.10, 0.64, 0, stroke=FAINT, stroke_width=1.2)
    s.arrow(8.31, 3.10, 0.64, 0, stroke=FAINT, stroke_width=1.2)
    s.rect(1.05, 4.36, 11.23, 1.36, fill=PAPER, stroke=BURGUNDY,
           stroke_width=1.1, radius=0.08)
    label(s, 1.34, 4.61, 0.88, "Thesis", fill=PALE_RED, ink=BURGUNDY)
    s.text(1.33, 5.03, 10.66, 0.42,
           "Train the sender and receiver together so the information needed for classification survives the link—then compare it fairly with a strong conventional system.",
           size=16.5, font="serif", text_color=INK, align="center", valign="mid")
    s.text(2.33, 6.20, 8.67, 0.28,
           "A simulation-first capstone with an optional real-radio demonstration later.",
           size=11.4, italic=True, font="sans", text_color=MUTED, align="center")
    add_footer(s); slides.append(s)

    # 2 — the application problem
    s = SlideScene(2, "Motivation", "The AI is only as useful as the link feeding it",
                   ("Motivation",),
                   "Original proposal: Semantic Communication + AI · edge/IoT application framing")
    add_header(s)
    s.text(0.82, 1.46, 11.69, 0.42,
           "Imagine a remote camera sending images to a server that decides what it sees.",
           size=17.0, font="serif", text_color=NAVY, bold=True, align="center")
    s.rect(0.82, 2.12, 7.31, 3.77, fill=PAPER, stroke=LINE, radius=0.08)
    label(s, 1.08, 2.37, 1.25, "Real-world tension")
    problems = [
        ("The sensor is constrained", "Bandwidth, latency and energy are limited at the edge."),
        ("The wireless link is imperfect", "Noise can corrupt a short message or make decoding fail."),
        ("The receiver has a narrow goal", "For classification, the final question is often simply: what is in the image?"),
    ]
    for i, (heading, body) in enumerate(problems):
        y = 2.96 + i * 0.90
        s.circle(1.11, y + 0.03, 0.30, 0.30, fill=BURGUNDY, stroke=None)
        s.text(1.11, y + 0.10, 0.30, 0.12, str(i + 1), size=7.8, bold=True,
               font="sans", text_color=WHITE, align="center")
        s.text(1.58, y, 2.27, 0.24, heading, size=11.5, bold=True,
               font="serif", text_color=INK)
        s.text(3.87, y - 0.02, 3.86, 0.40, body, size=10.6,
               font="sans", text_color=MUTED)
    s.rect(8.48, 2.12, 4.04, 3.77, fill=PALE_RED, stroke=None, radius=0.08)
    label(s, 8.77, 2.37, 1.18, "Design question", fill=WHITE, ink=BURGUNDY)
    s.text(8.83, 3.05, 3.34, 1.22,
           "Why spend the entire link budget reconstructing every pixel if the receiver ultimately needs a task decision?",
           size=16.4, bold=True, font="serif", text_color=INK, align="center")
    s.text(8.91, 4.72, 3.18, 0.54,
           "The project tests this question; it does not assume the learned method must win.",
           size=10.5, italic=True, font="sans", text_color=BURGUNDY, align="center")
    s.rect(1.75, 6.21, 9.83, 0.46, fill=PALE_AMBER, stroke=None, radius=0.06)
    s.text(1.99, 6.32, 9.35, 0.19,
           "Goal: preserve classification performance when communication resources are scarce and the link becomes noisy.",
           size=11.2, font="serif", bold=True, text_color=AMBER, align="center")
    add_footer(s); slides.append(s)

    # 3 — intuitive idea before the formal comparison
    s = SlideScene(3, "Core idea", "From bit-perfect delivery to task success",
                   ("Motivation", "Objectives"),
                   "Conceptual behavior only—not measured project results · formal three-arm experiment begins on slide 6")
    add_header(s)
    lanes = [
        (1.55, "CONVENTIONAL COMMUNICATION", "Compress the image → protect the bits → reconstruct the file → classify it", NAVY, PALE_NAVY),
        (3.38, "TASK-ORIENTED COMMUNICATION", "Learn a representation → send it through noise → recover what the task needs", BURGUNDY, PALE_RED),
    ]
    for y, heading, body, accent, pale in lanes:
        s.rect(0.82, y, 11.68, 1.36, fill=PAPER, stroke=LINE, radius=0.07)
        s.rect(0.82, y, 0.07, 1.36, fill=accent, stroke=None, radius=0)
        s.text(1.13, y + 0.23, 2.42, 0.25, heading, size=9.0, bold=True,
               font="sans", text_color=accent)
        s.text(3.48, y + 0.21, 8.56, 0.30, body, size=12.4,
               font="serif", text_color=INK)
        s.text(3.48, y + 0.78, 8.42, 0.27,
               "Primary objective: " + ("recover source bits and pixels" if accent == NAVY else "retain information useful for classification"),
               size=9.7, italic=True, font="sans", text_color=MUTED)
    s.rect(0.82, 5.20, 7.48, 1.23, fill=PALE_GREEN, stroke=None, radius=0.07)
    label(s, 1.08, 5.44, 1.42, "Expected signature", fill=WHITE, ink=GREEN)
    s.text(1.08, 5.91, 6.88, 0.26,
           "As the link worsens, a separated chain may hit a decoding cliff; a learned joint system may degrade more gradually.",
           size=11.5, font="serif", text_color=INK, align="center")
    s.rect(8.64, 5.20, 3.86, 1.23, fill=PALE_AMBER, stroke=None, radius=0.07)
    s.text(8.90, 5.45, 3.34, 0.20, "IMPORTANT CONTROL", size=8.2, bold=True,
           font="sans", text_color=AMBER, align="center")
    s.text(8.94, 5.84, 3.26, 0.36,
           "A digital system can also send learned features—so the experiment includes that third arm.",
           size=9.8, font="sans", text_color=INK, align="center")
    s.text(2.21, 6.68, 8.91, 0.20,
           "Research question: does the task-oriented joint design help after bandwidth, noise and comparison strength are matched fairly?",
           size=10.3, bold=True, font="serif", text_color=BURGUNDY, align="center")
    add_footer(s); slides.append(s)

    # 4 — literature synthesis
    s = SlideScene(4, "Problem survey", "How prior research leads to this project",
                   ("Problem Survey", "Subject Knowledge"),
                   "30-source synthesis in docs/literature-review.md · references [1]–[23]")
    add_header(s)
    boxes = [
        (0.72, 1.52, "Finite blocklength", "Rate backs off from capacity; source and channel dispersions interact.", "Need measured practical short-packet loss.", "[1–4]", NAVY),
        (6.83, 1.52, "Learned compression", "End-to-end transforms and entropy priors improve rate–distortion.", "Reconstruction optimization is not task optimization.", "[11–17]", BURGUNDY),
        (0.72, 3.65, "DeepJSCC", "Continuous learned mappings exhibit graceful degradation under mismatch.", "Most image studies foreground reconstruction quality.", "[5–10,24]", GREEN),
        (6.83, 3.65, "Task-oriented communication", "Learned representations can preserve remote inference under link constraints.", "Task awareness must be separated from joint coding.", "[18–23]", AMBER),
    ]
    for x, y, heading, premise, gap, refs, accent in boxes:
        s.rect(x, y, 5.78, 1.82, fill=PAPER, stroke=LINE, radius=0.08)
        s.rect(x, y, 0.07, 1.82, fill=accent, stroke=None, radius=0)
        s.text(x + 0.25, y + 0.22, 4.68, 0.30, heading, size=15.2, bold=True,
               font="serif", text_color=INK)
        s.text(x + 4.94, y + 0.24, 0.56, 0.22, refs, size=8.5, bold=True,
               font="mono", text_color=accent, align="right")
        s.text(x + 0.25, y + 0.69, 5.25, 0.42, premise, size=10.7,
               font="sans", text_color=MUTED)
        s.text(x + 0.25, y + 1.27, 1.20, 0.18, "OPEN ISSUE", size=7.5,
               bold=True, font="sans", text_color=accent)
        s.text(x + 1.30, y + 1.22, 4.18, 0.34, gap, size=9.4,
               italic=True, font="sans", text_color=INK)
    s.rect(1.24, 5.90, 10.84, 0.66, fill=PALE_NAVY, stroke=None, radius=0.07)
    s.text(1.48, 6.08, 10.36, 0.25,
           "Gap carried forward: three-way attribution at equal channel uses, with a strong adaptive classical comparator and paired failures retained.",
           size=11.3, font="serif", bold=True, text_color=NAVY, align="center")
    add_footer(s); slides.append(s)

    # 5 — objectives/hypotheses
    s = SlideScene(5, "Objectives & hypotheses", "What we will build—and what counts as success",
                   ("Objectives", "Hypothesis"),
                   "Exact decision protocol: SPEC §2 and ER-10 · H1 is the sole confirmatory primary")
    add_header(s)
    s.rect(0.70, 1.48, 4.50, 5.22, fill=PAPER, stroke=LINE, radius=0.08)
    label(s, 0.95, 1.72, 1.03, "Objectives")
    objectives = [
        "Build a reproducible normalized-AWGN DJSCC pipeline.",
        "Implement a standards-derived, non-strawman JPEG 2000 + NR LDPC baseline.",
        "Tune classical quality, code rate and modulation on validation only.",
        "Add the task-aware digital feature control for attribution.",
        "Evaluate paired image outcomes once on the sealed test split.",
        "Report positive, null or negative outcomes under one frozen protocol.",
    ]
    for i, item in enumerate(objectives, 1):
        y = 2.25 + (i - 1) * 0.68
        s.text(0.98, y, 0.33, 0.27, f"{i}.", size=10, bold=True, font="mono",
               text_color=BURGUNDY, align="right")
        s.text(1.45, y - 0.02, 3.35, 0.46, item, size=10.6, font="sans", text_color=INK)
    s.rect(5.55, 1.48, 7.08, 2.64, fill=PALE_RED, stroke=None, radius=0.08)
    label(s, 5.84, 1.74, 1.18, "H1 · primary", fill=WHITE, ink=BURGUNDY)
    s.text(5.84, 2.20, 6.48, 0.45,
           "Low-SNR separation: learned − classical_adaptive",
           size=15.2, bold=True, font="serif", text_color=INK)
    s.text(5.84, 2.79, 6.35, 0.80,
           "A point qualifies when √N·Δ̂(s)/σ̂(s) > 1.96. H1 is supported only if the longest run R_obs contains at least three consecutive qualifying points at or below the training SNR and the calibrated run p-value ≤ 0.05.",
           size=10.8, font="sans", text_color=MUTED)
    s.text(5.84, 3.67, 6.34, 0.20,
           "Effect size of record: mean paired accuracy difference over the full low-SNR region.",
           size=8.8, italic=True, font="sans", text_color=BURGUNDY)
    hypotheses = [
        ("H2", "Graceful vs cliff", "Fixed classical curve drops ≥30 pp while learned drops ≤15 pp over the frozen window."),
        ("H3", "Convergence", "Paired gap trends toward zero across SNR; crossover is reported, not required."),
        ("H4", "Attribution", "Credit joint coding only if DJSCC also exceeds the task-aware digital control under the H1 rule."),
    ]
    for i, (h, title, body) in enumerate(hypotheses):
        y = 4.37 + i * 0.72
        s.rect(5.55, y, 7.08, 0.58, fill=PAPER, stroke=LINE, radius=0.05)
        s.text(5.77, y + 0.16, 0.45, 0.18, h, size=9.3, bold=True,
               font="mono", text_color=NAVY)
        s.text(6.25, y + 0.12, 1.56, 0.20, title, size=10.2, bold=True,
               font="serif", text_color=INK)
        s.text(7.83, y + 0.10, 4.48, 0.30, body, size=8.9,
               font="sans", text_color=MUTED)
    add_footer(s); slides.append(s)

    # 6 — architecture
    s = SlideScene(6, "Architecture", "The experiment in one picture",
                   ("Subject Knowledge", "Objectives"),
                   "SPEC SR-4, SR-17, BR-4, ER-9 · src/data/test_access.py is the sole release boundary")
    add_header(s)
    s.rect(0.72, 1.48, 1.42, 4.75, fill=PALE_NAVY, stroke=None, radius=0.07)
    s.text(0.93, 1.78, 1.00, 0.25, "IMAGENETTE-160", size=9.0, bold=True,
           font="sans", text_color=NAVY, align="center")
    s.text(0.91, 2.26, 1.04, 0.70, "canonical\n160 × 160 RGB\npreprocessing", size=10.0,
           font="sans", text_color=INK, align="center")
    s.text(0.91, 3.40, 1.04, 0.58, "stable sample ID\n+ split manifest", size=9.2,
           font="mono", text_color=MUTED, align="center")
    s.arrow(2.18, 3.80, 0.42, 0, stroke=FAINT, stroke_width=1.2)
    arm_specs = [
        (2.18, "CLASSICAL IMAGE LINK", "JPEG 2000\npacketise + CRC\nNR LDPC + modulation", "decode\nreconstruct\nfrozen classifier", NAVY, PALE_NAVY),
        (3.58, "DIGITAL TASK-FEATURE LINK", "learned features\nquantize\nsame LDPC + modulation", "decode\nfeature task head", BURGUNDY, PALE_RED),
        (4.98, "LEARNED JOINT LINK", "neural encoder\npower normalize", "dual-head decoder\nreconstruction + class", GREEN, PALE_GREEN),
    ]
    for y, heading, left, right, accent, pale in arm_specs:
        s.rect(2.66, y, 3.00, 1.10, fill=PAPER, stroke=LINE, radius=0.06)
        s.text(2.88, y + 0.12, 2.56, 0.18, heading, size=8.5, bold=True,
               font="sans", text_color=accent, align="center")
        s.text(2.88, y + 0.42, 2.56, 0.50, left, size=9.1,
               font="sans", text_color=INK, align="center")
        s.arrow(5.73, y + 0.55, 0.42, 0, stroke=accent, stroke_width=1.2)
        s.rect(6.24, y, 1.86, 1.10, fill=pale, stroke=None, radius=0.06)
        s.text(6.44, y + 0.36, 1.46, 0.34, "AWGN\nkeyed noise", size=10.0,
               bold=True, font="math", text_color=accent, align="center")
        s.arrow(8.18, y + 0.55, 0.42, 0, stroke=accent, stroke_width=1.2)
        s.rect(8.69, y, 3.00, 1.10, fill=PAPER, stroke=LINE, radius=0.06)
        s.text(8.92, y + 0.35, 2.54, 0.48, right, size=9.2,
               font="sans", text_color=INK, align="center")
    s.rect(2.66, 6.25, 9.03, 0.38, fill=PALE_AMBER, stroke=None, radius=0.06)
    s.text(2.88, 6.34, 8.61, 0.16,
           "Per-image paired row: verdict + exact k + measured Eₛ + noise identity + task outcome",
           size=10.2, bold=True, font="mono", text_color=AMBER, align="center")
    s.rect(0.72, 6.38, 1.42, 0.28, fill=PALE_RED, stroke=None, radius=0.04)
    s.text(0.79, 6.44, 1.28, 0.13, "TEST SEALED → G-12", size=7.4, bold=True,
           font="sans", text_color=BURGUNDY, align="center")
    add_footer(s); slides.append(s)

    # 7 — fairness
    s = SlideScene(7, "Methodology", "How we keep the comparison fair",
                   ("Subject Knowledge", "Objectives"),
                   "BR-4, BR-10/11, ER-3/4/9/10 · W4 integration adjudication")
    add_header(s)
    controls = [
        ("Equal resource", "Same complex-symbol budget k and Eₛ/N₀ definition; realized symbol energy is logged.", NAVY),
        ("Strong comparator", "JPEG 2000 quality, LDPC rate and BPSK/QPSK/16-QAM tune per SNR on validation.", NAVY),
        ("Exact accounting", "Container, CRC, code-block and filler bytes count; decode failures remain in the denominator.", BURGUNDY),
        ("Measured BLER", "Every required physical identity is characterized; missing evidence is never treated as zero BLER.", BURGUNDY),
        ("Outage policy", "Constant-class fallback is measured from validation counts—100/1,000—not assumed as 1/10.", GREEN),
        ("Classifier control", "Clean and artifact-finetuned classifier passes separate codec shift from link failure.", GREEN),
        ("Attribution arm", "Quantized learned features traverse the same digital chain, isolating task awareness from joint coding.", AMBER),
        ("One test campaign", "All choices freeze before G-12; paired inference retains complete image × system trajectories.", AMBER),
    ]
    for i, (heading, body, accent) in enumerate(controls):
        col, row = i % 2, i // 2
        x = 0.72 + col * 6.05
        y = 1.50 + row * 1.28
        s.rect(x, y, 5.78, 1.02, fill=PAPER, stroke=LINE, radius=0.06)
        s.rect(x, y, 0.07, 1.02, fill=accent, stroke=None, radius=0)
        s.text(x + 0.24, y + 0.17, 1.47, 0.25, heading, size=11.0,
               bold=True, font="serif", text_color=INK)
        s.text(x + 1.72, y + 0.13, 3.80, 0.47, body, size=9.6,
               font="sans", text_color=MUTED)
    s.rect(1.55, 6.62, 10.23, 0.28, fill=PALE_RED, stroke=None, radius=0.04)
    s.text(1.76, 6.68, 9.81, 0.14,
           "Fail-closed rule: uncharacterized, infeasible and failed cases are explicit outcomes—not silently dropped rows.",
           size=8.6, bold=True, font="sans", text_color=BURGUNDY, align="center")
    add_footer(s); slides.append(s)

    # 8 — evidence
    s = SlideScene(8, "Evidence", "Engineering and gate evidence—not learned-vs-classical results",
                   ("Subject Knowledge",),
                   "Authenticated records: G-1, G-2, G-7, W4 and G8_C · test split access = 0")
    add_header(s)
    s.rect(0.72, 1.43, 11.88, 0.40, fill=PALE_RED, stroke=None, radius=0.04)
    s.text(0.96, 1.53, 11.40, 0.16,
           "NO HEADLINE COMPARISON EXISTS YET · bounded smoke is integration evidence, not scientific outcome evidence",
           size=8.7, bold=True, font="sans", text_color=BURGUNDY, align="center")
    evidence = [
        ("G-1 · classifier", "89.8%", "898/1,000 validation top-1\n88% preregistered floor · PASS", NAVY),
        ("G-7 · feasibility", "1.64 M", "DJSCC parameters · batch 32\n48.68 s/epoch · 1.004 GiB VRAM", GREEN),
        ("G-2 · digital PHY", "≤ 0.0037 dB", "observed waterfall displacement\nall three modulations within 0.5 dB · PASS", BURGUNDY),
        ("W4 · integration", "end-to-end", "JPEG 2000 → LDPC → AWGN → decode\nrecords + verifier · bounded PASS", AMBER),
    ]
    for i, (heading, metric, body, accent) in enumerate(evidence):
        x = 0.72 + i * 3.00
        s.rect(x, 2.10, 2.72, 2.62, fill=PAPER, stroke=LINE, radius=0.07)
        s.rect(x, 2.10, 2.72, 0.08, fill=accent, stroke=None, radius=0)
        s.text(x + 0.20, 2.42, 2.32, 0.24, heading, size=9.8, bold=True,
               font="sans", text_color=accent, align="center")
        s.text(x + 0.18, 3.02, 2.36, 0.45, metric, size=22, bold=True,
               font="serif", text_color=INK, align="center")
        s.text(x + 0.20, 3.78, 2.32, 0.56, body, size=9.1,
               font="sans", text_color=MUTED, align="center")
    s.rect(0.72, 5.10, 11.88, 1.22, fill=PALE_NAVY, stroke=None, radius=0.08)
    label(s, 0.98, 5.35, 1.18, "G8_C · current", fill=WHITE, ink=NAVY)
    s.text(2.38, 5.30, 2.46, 0.42, "paused", size=22, bold=True,
           font="serif", text_color=NAVY, align="center")
    s.text(4.92, 5.30, 4.28, 0.47,
           "authenticated full-strength BLER units\nsource epoch 2 registered · resumable campaign paused",
           size=10.4, font="sans", text_color=INK, align="center")
    s.text(9.39, 5.27, 2.83, 0.54,
           "BLER table: not frozen\nselection / inference / training / test: 0",
           size=9.8, bold=True, font="sans", text_color=BURGUNDY, align="center")
    s.text(1.20, 6.58, 10.95, 0.22,
           "Interpretation: the reference task, computational budget and digital chain are verified; the comparison itself remains ahead.",
           size=10.1, italic=True, font="serif", text_color=MUTED, align="center")
    add_footer(s); slides.append(s)

    # 9 — plan
    s = SlideScene(9, "Time plan", "Gate-ordered path from characterization to final report",
                   ("Time Plan",),
                   "docs/gantt-plan.md · dates from current 2026–27 capstone circular")
    add_header(s)
    s.text(0.76, 1.43, 11.79, 0.36,
           "Later work is shown as planned—not backfilled as complete for Review 1.",
           size=12.2, italic=True, font="serif", text_color=MUTED, align="center")
    stages = [
        ("G8_C", "BLER table", "CURRENT", 1.20, NAVY),
        ("G-8", "ratios + baseline", "NEXT", 1.28, NAVY),
        ("W5–W8", "DJSCC train + λ", "PLANNED", 1.52, GREEN),
        ("G-11", "digital control", "PLANNED", 1.38, BURGUNDY),
        ("G-12", "freeze + test", "SEALED", 1.38, AMBER),
        ("G-5", "demo decision", "PLANNED", 1.30, AMBER),
        ("W13–W17", "report + demo", "PLANNED", 1.60, NAVY),
    ]
    x = 0.72
    for i, (name, desc, status, w, accent) in enumerate(stages):
        pale = PALE_NAVY if accent == NAVY else PALE_GREEN if accent == GREEN else PALE_RED if accent == BURGUNDY else PALE_AMBER
        s.rect(x, 2.28, w, 1.30, fill=pale, stroke=None, radius=0.05)
        s.text(x + 0.10, 2.51, w - 0.20, 0.22, name, size=11.8, bold=True,
               font="serif", text_color=accent, align="center")
        s.text(x + 0.10, 2.92, w - 0.20, 0.28, desc, size=8.7,
               font="sans", text_color=INK, align="center")
        s.text(x + 0.10, 3.31, w - 0.20, 0.14, status, size=6.8, bold=True,
               font="mono", text_color=accent, align="center")
        if i < len(stages) - 1:
            s.arrow(x + w + 0.03, 2.92, 0.22, 0, stroke=FAINT, stroke_width=1.0)
        x += w + 0.30
    s.line(0.87, 4.36, 11.60, 0, stroke=LINE, stroke_width=1.2)
    milestones = [
        (1.02, "18–22 AUG", "FIRST REVIEW", BURGUNDY),
        (4.03, "29 SEP–3 OCT", "SECOND REVIEW", NAVY),
        (7.35, "17–21 NOV", "FINAL REVIEW", NAVY),
        (10.45, "20 NOV", "REPORT DUE", BURGUNDY),
    ]
    for x, date, title, accent in milestones:
        s.circle(x, 4.22, 0.28, 0.28, fill=accent, stroke=None)
        s.text(x - 0.22, 4.76, 1.72, 0.24, date, size=9.2, bold=True,
               font="mono", text_color=accent, align="left")
        s.text(x - 0.22, 5.14, 1.80, 0.22, title, size=8.1, bold=True,
               font="sans", text_color=INK, align="left")
    s.rect(1.10, 5.93, 11.08, 0.64, fill=PAPER, stroke=LINE, radius=0.06)
    s.text(1.35, 6.11, 10.58, 0.24,
           "Critical path: G8_C → G-8 → training → λ calibration → ER-9/G-11 → validation rehearsal → G-12 → one test campaign → G-5 → report",
           size=10.2, bold=True, font="mono", text_color=NAVY, align="center")
    s.text(3.33, 6.70, 6.68, 0.18,
           "W16 is report contingency—not room for new experiment scope.",
           size=8.8, italic=True, font="serif", text_color=MUTED, align="center")
    add_footer(s); slides.append(s)

    # 10 — standards/deployment/decision
    s = SlideScene(10, "Scope & decision", "What this First Review should confirm",
                   ("Subject Knowledge", "Objectives"),
                   "PR-3 standards register · PR-9 deployment dossier · guide acknowledgement remains PENDING")
    add_header(s)
    card(s, 0.72, 1.52, 3.73, 3.74, "Project boundary",
         "One narrow task: classify an image after communication over a bandwidth-limited noisy link.\n\nThe sender runs an encoder; the receiver owns the task head. That deployment split is fixed.\n\nTier 1 uses a normalized simulated AWGN channel so the scientific comparison is controlled and reproducible.\n\nThe result may be positive, null or negative.",
         accent=NAVY, body_size=10.5, tag="research scope")
    card(s, 4.80, 1.52, 3.73, 3.74, "Deployment ladder",
         "TIER 1 · required\nSimulation-first, offline and reproducible.\n\nTIER 2 · stretch after G-5\nConducted SDR I/Q replay; prerecorded outcome expected.\n\nTIER 3 · stretch\nEdge placement only if latency and procurement gates pass.\n\nCandidate: HackRF One + attenuation chain + RTL-SDR; no purchase yet.",
         accent=GREEN, body_size=10.2, tag="implementation")
    card(s, 8.88, 1.52, 3.73, 3.74, "Scope summary",
         "Confirm that the project asks the right question and compares the systems fairly:\n\n• image-classification is the task\n• communication resources are matched\n• the classical baseline is strong\n• the digital feature arm supports attribution\n• simulation-first Tier 1 is sufficient\n\nGuide acknowledgement: PENDING",
         accent=BURGUNDY, body_size=10.3, tag="panel ask", fill=PALE_RED)
    s.rect(1.04, 5.67, 11.24, 0.87, fill=PAPER, stroke=LINE, radius=0.07)
    s.text(1.31, 5.90, 2.30, 0.24, "TECHNICAL BOUNDARY", size=8.5, bold=True,
           font="mono", text_color=BURGUNDY, align="center")
    s.text(3.58, 5.85, 8.34, 0.38,
           "OpenJPEG 2.5.4 + TS 38.212-derived LDPC/rate matching over abstract AWGN—not a full 5G NR link claim",
           size=10.6, font="serif", text_color=INK, align="center")
    s.text(2.22, 6.72, 8.92, 0.18,
           "No guaranteed win, measured energy saving or real-radio generalization; completion means executing the preregistered protocol.",
           size=9.2, italic=True, font="serif", text_color=MUTED, align="center")
    add_footer(s); slides.append(s)

    # 11 — appendix rubric map
    s = SlideScene(11, "Appendix A", "Explicit First Review rubric-to-slide map",
                   ("Rubric Map",),
                   "Six criteria × five sub-marks · 30 scaled to 10 · final human readiness remains evidence-based")
    add_header(s)
    rows = [
        ("Motivation", "2–3", "Short packets, fixed k, task accuracy, finite-blocklength costs"),
        ("Objectives", "3, 5–7, 10", "Three-arm build, fair protocol, attribution control, decision ask"),
        ("Hypothesis", "5", "Exact H1 run rule; H2–H4 preregistered; crossing not required"),
        ("Problem Survey", "2, 4, 12", "30-source synthesis across four literature families"),
        ("Subject Knowledge", "3–8, 10", "AWGN, LDPC, JPEG 2000, fairness, evidence and scope boundary"),
        ("Time Plan", "9", "Gate-ordered critical path and fixed review/report dates"),
    ]
    s.rect(0.82, 1.55, 11.70, 0.52, fill=NAVY, stroke=None, radius=0.04)
    s.text(1.05, 1.72, 2.04, 0.18, "CRITERION", size=8.2, bold=True,
           font="sans", text_color=WHITE)
    s.text(3.38, 1.72, 1.24, 0.18, "SLIDES", size=8.2, bold=True,
           font="sans", text_color=WHITE)
    s.text(4.72, 1.72, 7.30, 0.18, "EVIDENCE EXPOSED IN THE DECK", size=8.2,
           bold=True, font="sans", text_color=WHITE)
    for i, (criterion, slide_nums, evidence) in enumerate(rows):
        y = 2.12 + i * 0.71
        fill = PAPER if i % 2 == 0 else IVORY
        s.rect(0.82, y, 11.70, 0.64, fill=fill, stroke=LINE, stroke_width=0.5, radius=0)
        s.text(1.05, y + 0.17, 2.05, 0.22, criterion, size=11.0, bold=True,
               font="serif", text_color=INK)
        s.text(3.38, y + 0.17, 1.15, 0.20, slide_nums, size=9.6, bold=True,
               font="mono", text_color=BURGUNDY)
        s.text(4.72, y + 0.14, 7.30, 0.27, evidence, size=9.5,
               font="sans", text_color=MUTED)
    s.rect(1.21, 6.56, 10.93, 0.33, fill=PALE_AMBER, stroke=None, radius=0.04)
    s.text(1.40, 6.64, 10.55, 0.15,
           "Final PASS still requires the real PPT review, four-member rehearsal, guide acknowledgement and annotated review-1-basis tag.",
           size=8.2, bold=True, font="sans", text_color=AMBER, align="center")
    add_footer(s); slides.append(s)

    # 12 — references and provenance
    s = SlideScene(12, "Appendix B", "Selected references and evidence provenance",
                   ("Problem Survey", "Subject Knowledge"),
                   "Full 30-reference synthesis: docs/literature-review.md · all evidence verified from repository records")
    add_header(s)
    s.text(0.72, 1.48, 5.72, 0.25, "SELECTED LITERATURE", size=9.0, bold=True,
           font="mono", text_color=BURGUNDY)
    refs_left = [
        "[1] Shannon, A Mathematical Theory of Communication, 1948.",
        "[2] Gastpar, Rimoldi & Vetterli, To Code, or Not to Code, 2003.",
        "[3] Polyanskiy, Poor & Verdú, Finite Blocklength, 2010.",
        "[4] Kostina & Verdú, Lossy JSCC at Finite Blocklength, 2013.",
        "[5] Bourtsoulatze, Kurka & Gündüz, DeepJSCC, 2019.",
        "[6] Kurka & Gündüz, DeepJSCC-f, 2020.",
        "[7] Kurka & Gündüz, Bandwidth-Agile DeepJSCC, 2021.",
        "[11] Ballé, Laparra & Simoncelli, Learned Compression, 2017.",
        "[12] Ballé et al., Scale Hyperprior, 2018.",
        "[17] Blau & Michaeli, Perception–Distortion, 2018.",
        "[18] Jankowski, Gündüz & Mikolajczyk, Wireless Retrieval, 2021.",
        "[19] Shao, Mao & Zhang, Task-Oriented Edge Inference, 2021.",
        "[20] Xie et al., Deep Learning Enabled Semantic Communication, 2021.",
    ]
    for i, ref in enumerate(refs_left):
        s.text(0.74, 1.90 + i * 0.34, 5.76, 0.23, ref, size=8.1,
               font="serif", text_color=INK)
    s.line(6.66, 1.50, 0, 4.82, stroke=LINE, stroke_width=0.8)
    s.text(6.98, 1.48, 5.60, 0.25, "REPOSITORY EVIDENCE", size=9.0, bold=True,
           font="mono", text_color=NAVY)
    evidence_paths = [
        ("G-1", "results/reference_classifier/g1_adjudication.json", "89.8% validation; test sealed"),
        ("G-7", "results/profiling/g7_djscc_profile.json", "1.64M params; measured epoch/VRAM"),
        ("G-2", "results/baseline/g2/g2_adjudication.json", "golden vectors + BLER conformance"),
        ("W4", "results/baseline/w4/integration_adjudication.json", "bounded end-to-end integration"),
        ("G8_C", "results/baseline/g8/campaign_state.json", "authenticated characterization cursor"),
        ("Spec", "spec/SPEC.md", "normative hypotheses and requirements"),
        ("Plan", "docs/gantt-plan.md", "dates, gates and current status"),
        ("Deploy", "docs/deployment-dossier.md", "simulation-first + SDR stretch"),
    ]
    for i, (tag, path, desc) in enumerate(evidence_paths):
        y = 1.91 + i * 0.57
        s.text(7.00, y, 0.62, 0.20, tag, size=8.2, bold=True,
               font="mono", text_color=BURGUNDY)
        s.text(7.70, y - 0.02, 4.55, 0.20, path, size=7.2,
               font="mono", text_color=NAVY)
        s.text(7.70, y + 0.23, 4.55, 0.17, desc, size=7.5,
               font="sans", text_color=MUTED)
    s.rect(6.98, 6.56, 5.20, 0.34, fill=PALE_RED, stroke=None, radius=0.04)
    s.text(7.14, 6.65, 4.88, 0.14,
           "Refresh G8_C status from instructions/RESUME.md on submission day.",
           size=7.6, bold=True, font="sans", text_color=BURGUNDY, align="center")
    add_footer(s); slides.append(s)

    return slides


def minimal_header(scene: SlideScene, title: str) -> None:
    scene.text(0.62, 0.48, 12.05, 0.55, title, size=27, bold=False,
               font="sans", text_color=INK)
    scene.line(0.62, 1.18, 12.05, 0, stroke=LINE, stroke_width=0.8)


def minimal_footer(scene: SlideScene) -> None:
    scene.line(0.62, 7.02, 12.05, 0, stroke=LINE, stroke_width=0.6)
    criterion = ", ".join(scene.criteria)
    scene.text(0.62, 7.12, 3.35, 0.16, f"Rubric: {criterion}", size=7.0,
               font="sans", text_color=MUTED)
    if scene.citation:
        scene.text(3.45, 7.12, 8.55, 0.16, scene.citation, size=6.8,
                   font="sans", text_color=MUTED, align="right")
    scene.text(12.12, 7.11, 0.50, 0.16, f"{scene.number} / 12", size=7.0,
               font="sans", text_color=MUTED, align="right")


def minimal_bullet(scene: SlideScene, y: float, text_value: str, *, x: float = 0.92,
                   w: float = 11.25, size: float = 15, h: float = 0.50) -> None:
    scene.text(x, y, 0.28, 0.28, "–", size=size, font="sans", text_color=MUTED)
    scene.text(x + 0.36, y, w - 0.36, h, text_value, size=size,
               font="sans", text_color=MUTED)


def minimal_table_row(scene: SlideScene, y: float, widths: tuple[float, ...],
                      values: tuple[str, ...], *, x: float = 0.72, h: float = 0.62,
                      header: bool = False, sizes: tuple[float, ...] | None = None) -> None:
    fill = "F3F3F3" if header else WHITE
    scene.rect(x, y, sum(widths), h, fill=fill, stroke=LINE, stroke_width=0.6, radius=0)
    offset = x
    for i, (w, value) in enumerate(zip(widths, values)):
        if i:
            scene.line(offset, y, 0, h, stroke=LINE, stroke_width=0.5)
        scene.text(offset + 0.12, y + 0.12, w - 0.24, h - 0.20, value,
                   size=(sizes[i] if sizes else (9.2 if header else 10.0)),
                   bold=header, font="sans", text_color=INK if header else MUTED,
                   valign="mid")
        offset += w


def _build_pre_knowledge_transfer_scenes() -> list[SlideScene]:
    """Google developer-documentation style First Review narrative."""
    slides: list[SlideScene] = []

    # 1 — title
    s = SlideScene(1, "Title", "Semantic Communication + AI",
                   ("Motivation", "Objectives"),
                   "First Review · 18–22 August 2026")
    s.text(0.72, 1.05, 11.75, 1.15, "Semantic Communication for Edge AI",
           size=36, bold=True, font="sans", text_color=INK)
    s.text(0.74, 2.48, 11.30, 0.48,
           "Task-oriented deep joint source–channel coding for image classification",
           size=19, font="sans", text_color=MUTED)
    s.line(0.74, 3.18, 11.80, 0, stroke=LINE, stroke_width=0.8)
    s.text(0.74, 3.62, 10.90, 0.74,
           "This project compares three ways to transmit an image for remote classification over a bandwidth-constrained noisy channel.",
           size=17, font="sans", text_color=MUTED)
    s.text(0.74, 4.72, 5.20, 0.28, "Capstone First Review", size=14,
           font="sans", text_color=INK)
    s.text(0.74, 5.18, 5.20, 0.28, "Review window: 18–22 August 2026", size=13,
           font="sans", text_color=MUTED)
    s.text(0.74, 5.64, 5.20, 0.28, "Presentation slot: 15 minutes", size=13,
           font="sans", text_color=MUTED)
    s.text(8.55, 5.62, 3.70, 0.35, "Team / guide details to be added", size=11,
           font="sans", text_color=MUTED, align="right")
    minimal_footer(s); slides.append(s)

    # 2 — context and definition
    s = SlideScene(2, "Context", "Project context and basic definition",
                   ("Motivation",),
                   "Original proposal: Semantic Communication + AI")
    minimal_header(s, s.title)
    s.text(0.76, 1.52, 11.70, 0.68,
           "The application is remote image classification. An edge device observes an image and sends information across a limited, noisy link to a receiver that owns the classifier.",
           size=17, font="sans", text_color=MUTED)
    s.text(0.76, 2.50, 11.70, 0.38, "Semantic communication", size=20,
           bold=True, font="sans", text_color=INK)
    s.text(0.76, 3.02, 11.50, 0.88,
           "Semantic communication transmits the information required to preserve meaning or complete a downstream task. It does not require every original source bit to be reconstructed when the task does not need them.",
           size=17, font="sans", text_color=MUTED)
    minimal_bullet(s, 4.22, "Source: an image observed by an edge camera or sensor.", size=15)
    minimal_bullet(s, 4.92, "Channel constraint: fixed bandwidth and additive wireless noise.", size=15)
    minimal_bullet(s, 5.62, "Receiver task: predict the correct image class.", size=15)
    s.text(0.76, 6.40, 11.45, 0.34,
           "Tier 1 is simulation-first. SDR and edge-device deployment are optional later stages.",
           size=13, font="sans", text_color=INK)
    minimal_footer(s); slides.append(s)

    # 3 — problem statement
    s = SlideScene(3, "Problem", "Problem statement and research question",
                   ("Motivation", "Objectives"),
                   "SPEC §1–2 · completion is independent of outcome")
    minimal_header(s, s.title)
    s.text(0.76, 1.52, 11.70, 0.35, "Conventional approach", size=19,
           bold=True, font="sans", text_color=INK)
    s.text(0.76, 1.98, 11.55, 0.68,
           "A conventional link compresses an image, protects the bits with a channel code, reconstructs the image at the receiver, and then runs the classifier. Its interfaces optimize image fidelity and bit recovery.",
           size=15.5, font="sans", text_color=MUTED)
    s.text(0.76, 2.98, 11.70, 0.35, "Task-oriented approach", size=19,
           bold=True, font="sans", text_color=INK)
    s.text(0.76, 3.44, 11.55, 0.68,
           "A learned joint system trains the sender and receiver together through a differentiable channel. Its representation can be optimized directly for classification after channel noise.",
           size=15.5, font="sans", text_color=MUTED)
    s.rect(0.76, 4.52, 11.55, 1.30, fill=WHITE, stroke="777777", stroke_width=1.0, radius=0)
    s.text(1.02, 4.76, 2.10, 0.25, "Research question", size=12,
           bold=True, font="sans", text_color=INK)
    s.text(1.02, 5.17, 10.92, 0.38,
           "At the same bandwidth and channel conditions, can task-aware joint coding preserve classification accuracy better than a properly tuned separated system?",
           size=16, font="sans", text_color=INK)
    s.text(0.76, 6.24, 11.55, 0.48,
           "The project reports the measured result even if the learned system does not outperform the baseline.",
           size=13.5, font="sans", text_color=MUTED)
    minimal_footer(s); slides.append(s)

    # 4 — literature
    s = SlideScene(4, "Survey", "Related work and the gap addressed",
                   ("Problem Survey", "Subject Knowledge"),
                   "30-source synthesis: docs/literature-review.md")
    minimal_header(s, s.title)
    topics = [
        ("Finite-blocklength communication", "Short packets operate below asymptotic capacity and incur non-negligible coding overhead and error probability. [1–4]"),
        ("Learned image compression", "End-to-end transforms improve rate–distortion performance, but reconstruction quality is not the same as task accuracy. [11–17]"),
        ("Deep joint source–channel coding", "DeepJSCC studies continuous learned mappings and graceful degradation over noisy channels. Most image results emphasize reconstruction. [5–10,24]"),
        ("Task-oriented communication", "Learned representations can preserve inference performance, but task awareness and joint coding must be separated experimentally. [18–23]"),
    ]
    for i, (heading, body) in enumerate(topics):
        y = 1.48 + i * 1.16
        s.text(0.78, y, 3.40, 0.30, heading, size=14.5, bold=True,
               font="sans", text_color=INK)
        s.text(4.17, y - 0.02, 8.03, 0.60, body, size=13.0,
               font="sans", text_color=MUTED)
        s.line(0.78, y + 0.82, 11.42, 0, stroke=LINE, stroke_width=0.5)
    s.text(0.78, 6.25, 2.20, 0.28, "Gap addressed", size=14.5,
           bold=True, font="sans", text_color=INK)
    s.text(3.00, 6.20, 9.20, 0.46,
           "A three-way comparison at equal channel uses: adaptive classical image transmission, task-aware digital features, and learned joint coding.",
           size=13.0, font="sans", text_color=MUTED)
    minimal_footer(s); slides.append(s)

    # 5 — objectives and completion
    s = SlideScene(5, "Objectives", "Project objectives and completion criterion",
                   ("Objectives", "Hypothesis"),
                   "Formal hypothesis rules are listed in Appendix A")
    minimal_header(s, s.title)
    objectives = [
        "Implement a normalized simulated AWGN channel and a task-oriented DJSCC system.",
        "Implement a JPEG 2000 + 5G NR LDPC baseline that is tuned on validation data.",
        "Implement a task-aware digital feature system using the same digital physical layer.",
        "Match channel uses, SNR definition, image identity and noise identity across systems.",
        "Freeze all choices before one paired evaluation on the sealed test split.",
        "Report positive, null or negative findings using the preregistered analysis.",
    ]
    for i, item in enumerate(objectives, 1):
        y = 1.48 + (i - 1) * 0.72
        s.text(0.84, y, 0.42, 0.28, f"{i}.", size=14, font="sans", text_color=MUTED)
        s.text(1.38, y, 10.85, 0.44, item, size=14.2, font="sans", text_color=MUTED)
    s.line(0.78, 5.94, 11.42, 0, stroke=LINE, stroke_width=0.6)
    s.text(0.78, 6.14, 2.30, 0.28, "Completion criterion", size=14.5,
           bold=True, font="sans", text_color=INK)
    s.text(3.08, 6.10, 9.05, 0.50,
           "Tier 1 is complete when all three systems are implemented, fairly matched, evaluated once, and reported. Completion does not require a positive result.",
           size=12.5, font="sans", text_color=MUTED)
    minimal_footer(s); slides.append(s)

    # 6 — systems
    s = SlideScene(6, "Design", "Systems under comparison",
                   ("Subject Knowledge", "Objectives"),
                   "Imagenette-160 headline dataset · normalized AWGN channel")
    minimal_header(s, s.title)
    widths = (2.50, 6.20, 3.00)
    minimal_table_row(s, 1.48, widths, ("System", "Sender and channel path", "Receiver output"), header=True)
    minimal_table_row(s, 2.10, widths,
                      ("Classical image link", "JPEG 2000 → packet framing → 5G NR LDPC → modulation → AWGN", "Reconstructed image → frozen classifier"), h=1.00)
    minimal_table_row(s, 3.10, widths,
                      ("Digital feature link", "Learned features → quantization → same LDPC and modulation → AWGN", "Recovered features → task head"), h=1.00)
    minimal_table_row(s, 4.10, widths,
                      ("Learned joint link", "Neural encoder → power normalization → AWGN", "Dual-head neural decoder → class"), h=1.00)
    s.text(0.78, 5.52, 2.20, 0.28, "Shared inputs", size=14.5, bold=True,
           font="sans", text_color=INK)
    s.text(3.00, 5.48, 9.15, 0.55,
           "Same image, dataset split, complex-symbol budget, SNR definition and keyed channel realization.",
           size=13.5, font="sans", text_color=MUTED)
    s.text(0.78, 6.28, 2.20, 0.28, "Data isolation", size=14.5, bold=True,
           font="sans", text_color=INK)
    s.text(3.00, 6.24, 9.15, 0.55,
           "Validation is used for selection. The test split remains sealed until gate G-12.",
           size=13.5, font="sans", text_color=MUTED)
    minimal_footer(s); slides.append(s)

    # 7 — controls
    s = SlideScene(7, "Method", "Experimental controls and fairness rules",
                   ("Subject Knowledge", "Objectives"),
                   "SPEC BR-4, BR-10/11, ER-3/9/10")
    minimal_header(s, s.title)
    widths = (3.05, 8.65)
    minimal_table_row(s, 1.46, widths, ("Control", "Rule"), header=True)
    rows = [
        ("Resource matching", "Compare systems at the same complex-symbol budget and SNR definition."),
        ("Classical baseline strength", "Tune JPEG 2000 quality, LDPC rate and BPSK/QPSK/16-QAM per SNR on validation."),
        ("Complete accounting", "Count container, CRC, segmentation and filler overhead. Keep decode failures in the denominator."),
        ("BLER evidence", "Measure every required physical-layer identity. Missing evidence is not interpreted as zero error."),
        ("Attribution", "Use the digital feature system to separate task-aware representation from joint coding."),
        ("Evaluation", "Use paired per-image outcomes and one sealed test campaign after all choices are frozen."),
    ]
    for i, row in enumerate(rows):
        minimal_table_row(s, 2.08 + i * 0.72, widths, row, h=0.72)
    minimal_footer(s); slides.append(s)

    # 8 — status
    s = SlideScene(8, "Status", "Engineering and gate evidence—not learned-vs-classical results",
                   ("Subject Knowledge",),
                   "Authenticated G-1, G-2, G-7, W4 and G8_C records · test access = 0")
    minimal_header(s, s.title)
    widths = (1.25, 1.35, 4.40, 4.70)
    minimal_table_row(s, 1.46, widths, ("Gate", "Status", "Recorded evidence", "Meaning"), header=True)
    status_rows = [
        ("G-1", "PASS", "898/1,000 = 89.8% validation accuracy", "The reference classification task is viable."),
        ("G-7", "PASS", "1.64 M parameters; 48.68 s per epoch; 1.004 GiB VRAM", "The planned DJSCC model fits the available GPU."),
        ("G-2", "PASS", "Golden vectors and BLER waterfall checks within 0.5 dB", "The digital physical layer passes conformance checks."),
        ("W4", "PASS", "Bounded JPEG 2000 → LDPC → AWGN → decode run", "The classical pipeline and evidence path work end to end."),
        ("G8_C", "PAUSED", "Durable full-strength characterization checkpoint", "The full BLER table and operating-point selection are not complete."),
    ]
    for i, row in enumerate(status_rows):
        minimal_table_row(s, 2.08 + i * 0.80, widths, row, h=0.80,
                          sizes=(10.2, 9.5, 9.8, 9.8))
    s.text(0.78, 6.40, 11.45, 0.34,
           "No learned-versus-classical headline result exists yet. Bounded integration evidence is not presented as a scientific comparison.",
           size=12.5, bold=True, font="sans", text_color=INK)
    minimal_footer(s); slides.append(s)

    # 9 — schedule
    s = SlideScene(9, "Schedule", "Remaining work and review dates",
                   ("Time Plan",),
                   "docs/gantt-plan.md · unfinished stages are shown as unfinished")
    minimal_header(s, s.title)
    widths = (2.05, 1.55, 4.90, 3.20)
    minimal_table_row(s, 1.44, widths, ("Period / gate", "State", "Work", "Output"), header=True)
    schedule_rows = [
        ("Now · G8_C", "Active", "Complete BLER characterization and freeze the table", "Inputs for G-8"),
        ("G-8", "Planned", "Select operating ratios and classical configurations", "Frozen baseline choices"),
        ("W5–W8", "Planned", "Train DJSCC; calibrate λ; implement digital feature control", "Frozen learned systems"),
        ("W9–W11", "Planned", "Validation rehearsal; freeze; one test campaign", "Hypothesis decisions"),
        ("W13–W17", "Planned", "Demo, report, poster and final review preparation", "Final package"),
    ]
    for i, row in enumerate(schedule_rows):
        minimal_table_row(s, 2.06 + i * 0.78, widths, row, h=0.78,
                          sizes=(10.0, 9.5, 9.8, 9.8))
    s.text(0.78, 6.20, 11.45, 0.36,
           "First Review: 18–22 Aug · Second Review: 29 Sep–3 Oct · Final Review: 17–21 Nov · Report due: 20 Nov",
           size=12.4, font="sans", text_color=INK)
    s.text(0.78, 6.58, 11.45, 0.28,
           "W16 is reserved for report contingency. It is not available for new experiment scope.",
           size=11.3, font="sans", text_color=MUTED)
    minimal_footer(s); slides.append(s)

    # 10 — scope and ask
    s = SlideScene(10, "Review ask", "Scope, deliverables, and decisions requested",
                   ("Objectives", "Subject Knowledge"),
                   "PR-3 standards register · PR-9 deployment dossier · guide acknowledgement PENDING")
    minimal_header(s, s.title)
    sections = [
        ("Required Tier 1 scope", "Simulation-based comparison of the three systems on image classification over normalized AWGN, with one sealed test campaign and paired analysis."),
        ("Optional later scope", "Conducted SDR replay and Raspberry Pi deployment are stretch goals after gate G-5. No hardware has been purchased."),
        ("Technical claim boundary", "OpenJPEG 2.5.4 and a TS 38.212-derived LDPC/rate-matching chain over abstract AWGN. This is not a full 5G NR link claim."),
        ("Scope summary", "Task boundary, three-system comparison, fairness controls and simulation-first Tier 1."),
    ]
    for i, (heading, body) in enumerate(sections):
        y = 1.48 + i * 1.18
        s.text(0.78, y, 3.00, 0.30, heading, size=15, bold=True,
               font="sans", text_color=INK)
        s.text(3.80, y - 0.02, 8.38, 0.62, body, size=13.2,
               font="sans", text_color=MUTED)
        s.line(0.78, y + 0.82, 11.40, 0, stroke=LINE, stroke_width=0.5)
    s.text(0.78, 6.24, 3.00, 0.30, "Outstanding human item", size=15,
           bold=True, font="sans", text_color=INK)
    s.text(3.80, 6.20, 8.38, 0.50,
           "The guide's dated acknowledgement of the simulation-first, no-required-hardware path is still pending.",
           size=13.2, font="sans", text_color=MUTED)
    minimal_footer(s); slides.append(s)

    # 11 — formal analysis appendix
    s = SlideScene(11, "Appendix A", "Hypotheses and analysis plan",
                   ("Hypothesis", "Subject Knowledge"),
                   "Exact formal protocol: SPEC §2 and ER-10")
    minimal_header(s, s.title)
    s.text(0.78, 1.44, 11.50, 0.30, "H1 — Low-SNR separation (confirmatory primary)",
           size=15.5, bold=True, font="sans", text_color=INK)
    s.text(0.78, 1.88, 11.50, 1.02,
           "Comparator: learned versus classical_adaptive at the headline ratio. A grid point qualifies when the studentized paired mean √N·Δ̂(s)/σ̂(s) exceeds 1.96. Let R_obs be the longest run of consecutive qualifying points at or below the training SNR. H1 is supported only if R_obs ≥ 3 and the calibrated run p-value ≤ 0.05.",
           size=13.3, font="sans", text_color=MUTED)
    formal_rows = [
        ("H2 — Graceful versus cliff", "The fixed classical curve must show the preregistered large drop while the learned curve remains within its smaller-drop bound."),
        ("H3 — Convergence", "The paired learned-minus-classical gap must contract across SNR. A crossover is reported if observed but is not required."),
        ("H4 — Attribution", "Joint coding is credited only if the learned joint system also exceeds the task-aware digital feature control under the H1 procedure."),
    ]
    for i, (heading, body) in enumerate(formal_rows):
        y = 3.16 + i * 0.94
        s.text(0.78, y, 3.45, 0.30, heading, size=13.5, bold=True,
               font="sans", text_color=INK)
        s.text(4.23, y - 0.02, 7.95, 0.55, body, size=12.0,
               font="sans", text_color=MUTED)
        s.line(0.78, y + 0.66, 11.40, 0, stroke=LINE, stroke_width=0.5)
    s.text(0.78, 6.28, 11.45, 0.42,
           "All hypotheses use paired per-image outcomes. The test split is opened once after gate G-12.",
           size=12.5, font="sans", text_color=INK)
    minimal_footer(s); slides.append(s)

    # 12 — rubric map and references
    s = SlideScene(12, "Appendix B", "First Review rubric map and sources",
                   ("Problem Survey", "Time Plan", "Subject Knowledge"),
                   "Full literature review contains 30 synthesized references")
    minimal_header(s, s.title)
    widths = (2.10, 1.25, 4.70)
    minimal_table_row(s, 1.42, widths, ("Criterion", "Slides", "Evidence"), header=True)
    rubric_rows = [
        ("Motivation", "2–3", "Application context and problem statement"),
        ("Objectives", "3, 5–7, 10", "Systems, completion criterion and review ask"),
        ("Hypothesis", "5, 11", "Outcome-independent completion and formal rules"),
        ("Problem Survey", "4, 12", "30-source literature synthesis"),
        ("Subject Knowledge", "6–8, 10–11", "System design, controls, evidence and scope"),
        ("Time Plan", "9", "Gate-ordered schedule and fixed review dates"),
    ]
    for i, row in enumerate(rubric_rows):
        minimal_table_row(s, 2.04 + i * 0.54, widths, row, h=0.54,
                          sizes=(9.4, 9.0, 9.2))
    s.line(8.40, 1.42, 0, 4.80, stroke=LINE, stroke_width=0.7)
    s.text(8.70, 1.42, 3.45, 0.28, "Selected references", size=14,
           bold=True, font="sans", text_color=INK)
    refs = [
        "[1] Shannon, 1948.",
        "[3] Polyanskiy, Poor & Verdú, 2010.",
        "[4] Kostina & Verdú, 2013.",
        "[5] Bourtsoulatze, Kurka & Gündüz, 2019.",
        "[11] Ballé, Laparra & Simoncelli, 2017.",
        "[17] Blau & Michaeli, 2018.",
        "[18] Jankowski, Gündüz & Mikolajczyk, 2021.",
        "[20] Xie et al., 2021.",
    ]
    for i, ref in enumerate(refs):
        s.text(8.70, 1.92 + i * 0.46, 3.50, 0.24, ref, size=9.2,
               font="sans", text_color=MUTED)
    s.text(0.78, 5.66, 2.20, 0.28, "Evidence records", size=14,
           bold=True, font="sans", text_color=INK)
    s.text(0.78, 6.08, 7.30, 0.66,
           "G-1: results/reference_classifier/g1_adjudication.json\nG-7: results/profiling/g7_djscc_profile.json\nG-2 / W4 / G8_C: results/baseline/",
           size=9.3, font="mono", text_color=MUTED)
    s.text(8.70, 5.72, 3.50, 0.54,
           "Refresh G8_C status from instructions/RESUME.md on submission day.",
           size=9.3, font="sans", text_color=MUTED)
    minimal_footer(s); slides.append(s)

    return slides


def build_minimal_scenes() -> list[SlideScene]:
    """Build the newcomer-first Review 1 narrative.

    The ordering is deliberate: establish the human task, explain the ordinary
    solution, introduce the learned change, and only then present the formal
    experiment. This keeps the deck useful to a panel member who has not read
    the repository or the specification.
    """
    slides: list[SlideScene] = []

    # 1 — title and plain statement of scope
    s = SlideScene(1, "Title", "Semantic Communication for Edge AI",
                   ("Motivation", "Objectives"),
                   "First Review · 18–22 August 2026")
    s.text(0.74, 1.00, 11.80, 1.20,
           "Semantic Communication + AI",
           size=34, bold=True, font="sans", text_color=INK)
    s.text(0.76, 2.42, 11.35, 0.50,
           "Learning what to send for a remote image-classification task",
           size=19, font="sans", text_color=MUTED)
    s.line(0.76, 3.18, 11.78, 0, stroke=LINE, stroke_width=0.9)
    s.text(0.76, 3.62, 11.30, 0.84,
           "An edge camera sees an image. A server across a limited, noisy link must recognize it. We compare a normal communication system with systems that learn what information to send.",
           size=17, font="sans", text_color=MUTED)
    s.text(0.76, 4.82, 5.30, 0.30, "Capstone First Review", size=14,
           bold=True, font="sans", text_color=INK)
    s.text(0.76, 5.30, 5.80, 0.30, "Review window: 18–22 August 2026", size=13,
           font="sans", text_color=MUTED)
    s.text(0.76, 5.74, 5.80, 0.30, "Presentation slot: 15 minutes", size=13,
           font="sans", text_color=MUTED)
    s.text(8.10, 5.72, 4.10, 0.36, "Team and guide details to be added", size=11,
           font="sans", text_color=MUTED, align="right")
    minimal_footer(s); slides.append(s)

    # 2 — the whole project before any technical machinery
    s = SlideScene(2, "Overview", "What problem are we solving?",
                   ("Motivation", "Objectives"),
                   "The task is remote image classification over a limited noisy link")
    minimal_header(s, s.title)
    steps = [
        ("1", "Observe", "An edge camera sees an image."),
        ("2", "Communicate", "The link has limited bandwidth and added noise."),
        ("3", "Classify", "A server decides what the image contains."),
        ("4", "Question", "Must the camera send every image detail?"),
    ]
    x_positions = (0.76, 3.84, 6.92, 10.00)
    for x, (number, heading, body) in zip(x_positions, steps):
        s.rect(x, 1.64, 2.48, 2.18, fill=WHITE, stroke=LINE,
               stroke_width=0.9, radius=0)
        s.text(x + 0.18, 1.86, 0.34, 0.30, number, size=14, bold=True,
               font="sans", text_color=INK)
        s.text(x + 0.18, 2.34, 2.10, 0.34, heading, size=18, bold=True,
               font="sans", text_color=INK)
        s.text(x + 0.18, 2.88, 2.08, 0.64, body, size=13.2,
               font="sans", text_color=MUTED)
    s.rect(0.76, 4.42, 11.72, 1.25, fill="F5F5F5", stroke=None, radius=0)
    s.text(1.02, 4.68, 3.15, 0.30, "Why this matters", size=16,
           bold=True, font="sans", text_color=INK)
    s.text(4.18, 4.64, 7.98, 0.56,
           "Sending a complete high-quality image can be expensive when the receiver only needs a correct task decision.",
           size=15, font="sans", text_color=MUTED)
    s.text(0.78, 6.14, 11.55, 0.48,
           "The project asks whether a learned system can preserve that decision using the same channel budget.",
           size=15, bold=True, font="sans", text_color=INK, align="center")
    minimal_footer(s); slides.append(s)

    # 3 — ordinary system, with each term introduced in context
    s = SlideScene(3, "Conventional approach", "How is this normally done?",
                   ("Motivation", "Subject Knowledge"),
                   "Headline baseline: JPEG 2000 + TS 38.212-derived LDPC over AWGN")
    minimal_header(s, s.title)
    path = [
        (0.70, 1.62, 1.55, "Image", "camera input"),
        (2.72, 1.62, 1.80, "Compress", "JPEG 2000"),
        (4.99, 1.62, 2.04, "Protect bits", "LDPC + modulation"),
        (7.50, 1.62, 1.74, "Add noise", "simulated channel"),
        (9.72, 1.62, 2.90, "AI classifier", "recovered image → ResNet-18"),
    ]
    for x, y, w, heading, body in path:
        s.rect(x, y, w, 1.34, fill=WHITE, stroke="777777", stroke_width=0.8, radius=0)
        s.text(x + 0.12, y + 0.27, w - 0.24, 0.30, heading, size=14,
               bold=True, font="sans", text_color=INK, align="center")
        s.text(x + 0.12, y + 0.78, w - 0.24, 0.28, body, size=11,
               font="sans", text_color=MUTED, align="center")
    for x in (2.34, 4.61, 7.12, 9.33):
        s.arrow(x, 2.28, 0.28, 0, stroke=MUTED, stroke_width=1.0)
    definitions = [
        ("Compression", "reduces the number of source bits."),
        ("Error correction", "adds controlled redundancy so damaged bits can be recovered."),
        ("Modulation", "maps protected bits to symbols sent through the channel."),
        ("AI classifier", "a frozen ResNet-18 maps the recovered image to class logits."),
    ]
    for i, (term, meaning) in enumerate(definitions):
        y = 3.48 + i * 0.58
        s.text(0.80, y, 2.12, 0.28, term, size=13.5, bold=True,
               font="sans", text_color=INK)
        s.text(2.98, y - 0.02, 9.26, 0.34, meaning, size=13.2,
               font="sans", text_color=MUTED)
    s.line(0.78, 5.98, 11.46, 0, stroke=LINE, stroke_width=0.6)
    s.text(0.80, 6.20, 11.42, 0.44,
           "Limitation: the communication stack is designed separately from the AI objective, even though task accuracy is the final metric.",
           size=14.2, bold=True, font="sans", text_color=INK)
    minimal_footer(s); slides.append(s)

    # 4 — learned change, including the class-label objection
    s = SlideScene(4, "Proposed approach", "What are we proposing?",
                   ("Motivation", "Objectives", "Subject Knowledge"),
                   "DJSCC = deep joint source–channel coding")
    minimal_header(s, s.title)
    learned_path = [
        (0.86, 1.62, 1.72, "Image", "sender input"),
        (3.10, 1.62, 2.14, "Learned sender", "decides what to send"),
        (5.76, 1.62, 1.86, "Noisy link", "limited bandwidth"),
        (8.14, 1.62, 2.30, "Learned receiver", "recovers task information"),
        (10.96, 1.62, 1.46, "Class", "final decision"),
    ]
    for x, y, w, heading, body in learned_path:
        s.rect(x, y, w, 1.34, fill=WHITE, stroke="777777", stroke_width=0.8, radius=0)
        s.text(x + 0.10, y + 0.27, w - 0.20, 0.30, heading, size=13.5,
               bold=True, font="sans", text_color=INK, align="center")
        s.text(x + 0.10, y + 0.76, w - 0.20, 0.34, body, size=10.7,
               font="sans", text_color=MUTED, align="center")
    for x in (2.68, 5.34, 7.72, 10.54):
        s.arrow(x, 2.28, 0.28, 0, stroke=MUTED, stroke_width=1.0)
    s.text(0.82, 3.48, 3.08, 0.34, "The basic idea", size=17,
           bold=True, font="sans", text_color=INK)
    s.text(3.92, 3.43, 8.22, 0.70,
           "Train the sender and receiver together using labeled images and simulated noise. The system can then learn which image information helps the receiver make the correct decision.",
           size=15, font="sans", text_color=MUTED)
    s.line(0.80, 4.36, 11.42, 0, stroke=LINE, stroke_width=0.6)
    s.text(0.82, 4.64, 3.08, 0.34, "Why not send the label?", size=17,
           bold=True, font="sans", text_color=INK)
    s.text(3.92, 4.58, 8.22, 0.88,
           "The project fixes the deployment split: the sender has an encoder, while the receiver owns the decoder and classifier. Sending a label would move the full task model to the sender and test a different system.",
           size=14.2, font="sans", text_color=MUTED)
    s.rect(0.80, 5.84, 11.42, 0.70, fill="F5F5F5", stroke=None, radius=0)
    s.text(1.02, 5.98, 10.98, 0.20,
           "Technical name: deep joint source–channel coding (DJSCC)",
           size=13.3, bold=True, font="sans", text_color=INK, align="center")
    s.text(1.02, 6.28, 10.98, 0.18,
           "AIML method: supervised end-to-end training · model and loss details follow in the methodology",
           size=10.8, font="sans", text_color=MUTED, align="center")
    minimal_footer(s); slides.append(s)

    # 5 — literature and gap before the project objectives
    s = SlideScene(5, "Problem survey", "Literature survey and research gap",
                   ("Problem Survey", "Subject Knowledge"),
                   "30-source synthesis: docs/literature-review.md")
    minimal_header(s, s.title)
    topics = [
        ("Short-packet communication", "Known", "Finite packets pay coding overhead and can fail abruptly.", "Measure the practical baseline instead of assuming ideal capacity. [1–4]"),
        ("Learned image compression", "Known", "Neural transforms learn compact image representations.", "Good reconstruction quality is not the same as good AI task accuracy. [11–17]"),
        ("DeepJSCC", "Known", "Neural encoders and decoders can learn continuous channel mappings.", "Most image studies emphasize reconstruction rather than classification. [5–10,24]"),
        ("Task-oriented edge AI", "Known", "Learned features can preserve inference under link limits.", "Representation learning and joint coding are often not isolated. [18–23]"),
    ]
    for i, (family, _, known, open_issue) in enumerate(topics):
        y = 1.46 + i * 1.12
        s.text(0.78, y, 2.58, 0.34, family, size=14.5, bold=True,
               font="sans", text_color=INK)
        s.text(3.42, y - 0.02, 4.10, 0.54, known, size=12.8,
               font="sans", text_color=MUTED)
        s.text(7.72, y - 0.02, 4.48, 0.58, open_issue, size=12.8,
               font="sans", text_color=MUTED)
        s.line(0.78, y + 0.80, 11.42, 0, stroke=LINE, stroke_width=0.5)
    s.text(0.78, 6.18, 2.58, 0.30, "Gap addressed", size=14.5, bold=True,
           font="sans", text_color=INK)
    s.text(3.42, 6.12, 8.78, 0.50,
           "A controlled three-way experiment that separates neural representation learning from joint channel coding at equal channel use.",
           size=13.4, font="sans", text_color=MUTED)
    minimal_footer(s); slides.append(s)

    # 6 — problem statement, research question and completion objectives
    s = SlideScene(6, "Objectives", "Problem statement, research question and objectives",
                   ("Motivation", "Objectives"),
                   "Objectives are stated as completion work, not promised outcomes")
    minimal_header(s, s.title)
    s.text(0.78, 1.44, 2.46, 0.30, "Problem statement", size=15.5,
           bold=True, font="sans", text_color=INK)
    s.text(3.26, 1.40, 8.88, 0.74,
           "An edge camera has a limited noisy link, while the server only needs a correct class decision. A conventional system tries to rebuild the image first. The project tests whether a learned system can protect the decision more directly.",
           size=13.4, font="sans", text_color=MUTED)
    s.line(0.78, 2.34, 11.42, 0, stroke=LINE, stroke_width=0.6)
    s.text(0.78, 2.62, 2.46, 0.30, "Research question", size=15.5,
           bold=True, font="sans", text_color=INK)
    s.text(3.26, 2.57, 8.88, 0.78,
           "At the same bandwidth and channel conditions, can the learned joint system preserve classification accuracy better than a tuned conventional system, and does joint training explain any remaining gain?",
           size=13.6, font="sans", text_color=INK)
    s.line(0.78, 3.52, 11.42, 0, stroke=LINE, stroke_width=0.6)
    s.text(0.78, 3.80, 2.46, 0.30, "Objectives", size=15.5,
           bold=True, font="sans", text_color=INK)
    objectives = [
        "Build and train the learned sender and receiver models.",
        "Build a digital feature system that separates task awareness from joint training.",
        "Match channel use, SNR, image identity and noise across all systems.",
        "Freeze all learned models and choices, then run one paired test campaign.",
    ]
    for i, item in enumerate(objectives, 1):
        y = 3.78 + (i - 1) * 0.64
        s.text(3.28, y, 0.36, 0.28, f"{i}.", size=12.6,
               font="sans", text_color=MUTED)
        s.text(3.72, y - 0.02, 8.42, 0.42, item, size=12.8,
               font="sans", text_color=MUTED)
    s.rect(0.78, 6.34, 11.42, 0.40, fill="F5F5F5", stroke=None, radius=0)
    s.text(0.98, 6.43, 11.02, 0.18,
           "Completion depends on executing the fixed protocol correctly, not on obtaining a favorable result.",
           size=12.2, bold=True, font="sans", text_color=INK, align="center")
    minimal_footer(s); slides.append(s)

    # 7 — proposed design and fair-comparison walkthrough
    s = SlideScene(7, "Methodology", "Proposed methodology",
                   ("Objectives", "Subject Knowledge"),
                   "Three systems · one shared task, resource budget and channel definition")
    minimal_header(s, s.title)
    s.rect(0.78, 1.42, 11.42, 0.62, fill="F5F5F5", stroke=None, radius=0)
    s.text(1.00, 1.60, 10.98, 0.26,
           "Imagenette-160 · same image and split · same complex-symbol budget · same SNR · same keyed noise",
           size=12.5, bold=True, font="sans", text_color=INK, align="center")
    widths = (2.36, 6.05, 3.29)
    minimal_table_row(s, 2.24, widths, ("System", "Transmission path", "Scientific role"),
                      header=True, sizes=(10.2, 10.2, 10.2))
    rows = [
        ("Classical + AI", "JPEG 2000 → LDPC → AWGN → recovered image → frozen ResNet-18", "Conventional reference"),
        ("AI feature control", "Shared neural encoder → quantization → same digital PHY → shared task head", "Isolates representation learning"),
        ("Neural DJSCC", "Residual encoder → normalization → AWGN → dual-head neural decoder", "Tests end-to-end neural coding"),
    ]
    for i, row in enumerate(rows):
        minimal_table_row(s, 2.86 + i * 0.94, widths, row, h=0.94,
                          sizes=(11.6, 11.0, 11.0))
    s.text(0.80, 5.86, 2.18, 0.28, "Evaluation controls", size=14.5,
           bold=True, font="sans", text_color=INK)
    s.text(2.98, 5.80, 9.18, 0.66,
           "Train with supervised CE + λMSE; tune on validation only; freeze before test access; retain failed transmissions; compare paired per-image AI outcomes.",
           size=12.7, font="sans", text_color=MUTED)
    s.rect(0.78, 6.48, 11.42, 0.28, fill="F5F5F5", stroke=None, radius=0)
    s.text(0.98, 6.53, 11.02, 0.14,
           "The test split remains sealed until gate G-12. No test result exists at First Review.",
           size=10.8, bold=True, font="sans", text_color=INK, align="center")
    minimal_footer(s); slides.append(s)

    # 8 — hypothesis and evaluation plan, now in the expected proposal position
    s = SlideScene(8, "Hypothesis", "Research hypotheses",
                   ("Hypothesis", "Subject Knowledge"),
                   "Primary metric: paired top-1 AI inference accuracy · exact protocol: SPEC §2")
    minimal_header(s, s.title)
    s.text(0.78, 1.42, 6.90, 0.30, "Primary hypothesis (H1)",
           size=15.5, bold=True, font="sans", text_color=INK)
    s.text(0.78, 1.86, 6.92, 0.90,
           "At difficult low-SNR conditions, the learned joint system will preserve classification accuracy better than the tuned conventional system across a consecutive region—not at only one isolated point.",
           size=14.0, font="sans", text_color=INK)
    s.text(0.78, 2.92, 6.90, 0.28, "Formal decision rule", size=13.2,
           bold=True, font="sans", text_color=INK)
    s.text(0.78, 3.28, 6.92, 1.34,
           "Comparator: learned versus classical_adaptive at the headline ratio. A point qualifies when √N·Δ̂(s)/σ̂(s) > 1.96. R_obs is the longest consecutive qualifying run at or below the training SNR. H1 is supported only if R_obs ≥ 3 and the calibrated run p-value ≤ 0.05.",
           size=11.8, font="sans", text_color=MUTED)
    s.text(0.78, 4.76, 6.92, 0.64,
           "Effect size: the mean paired accuracy difference over the full low-SNR region. A curve crossing is reported if observed, but is not required.",
           size=11.8, font="sans", text_color=MUTED)
    s.line(7.96, 1.42, 0, 4.68, stroke=LINE, stroke_width=0.7)
    secondary = [
        ("H2", "Does a fixed conventional configuration show a cliff while the learned system degrades more gradually?"),
        ("H3", "Does the learned–classical accuracy gap contract as SNR improves?"),
        ("H4", "Does the learned joint system also exceed the task-aware digital feature control?"),
    ]
    for i, (tag, question) in enumerate(secondary):
        y = 1.46 + i * 1.30
        s.text(8.24, y, 0.58, 0.28, tag, size=14.2, bold=True,
               font="sans", text_color=INK)
        s.text(8.90, y - 0.02, 3.24, 0.84, question, size=12.2,
               font="sans", text_color=MUTED)
        s.line(8.24, y + 0.98, 3.90, 0, stroke=LINE, stroke_width=0.45)
    s.rect(0.78, 5.68, 11.42, 0.84, fill="F5F5F5", stroke=None, radius=0)
    s.text(1.04, 5.92, 10.90, 0.38,
           "Outcome rule: support, no support or an adverse result are all valid findings. The protocol does not change after test access.",
           size=13.0, bold=True, font="sans", text_color=INK, align="center")
    minimal_footer(s); slides.append(s)

    # 9 — status in plain meaning, with no volatile campaign counts
    s = SlideScene(9, "Current status", "Preliminary work and feasibility",
                   ("Subject Knowledge",),
                   "Engineering and gate evidence only · no learned-versus-classical result")
    minimal_header(s, s.title)
    widths = (1.18, 3.62, 3.42, 3.48)
    minimal_table_row(s, 1.46, widths, ("Item", "Observed evidence", "What it tells us", "What it does not tell us"),
                      header=True, sizes=(10.0, 10.0, 10.0, 10.0))
    rows = [
        ("AI task", "ResNet-18 from scratch: 89.8% validation", "The classification task and clean reference model are viable.", "Validation only; test is sealed. (G-1)"),
        ("Learned model", "DJSCC: 1.64 M parameters", "The planned model fits the available RTX 4060.", "A feasibility profile, not training. (G-7)"),
        ("Digital link", "Golden vectors and BLER conformance pass", "The channel-coding implementation behaves as expected.", "Conformance evidence, not the full table. (G-2)"),
        ("Full pipeline", "Bounded conventional run completes", "Compression, coding, channel and records work together.", "Integration evidence, not a comparison. (W4)"),
        ("Current phase", "BLER characterization at a durable checkpoint", "The current measurement work can resume safely.", "Table, selection and training are unfinished. (G8_C)"),
    ]
    for i, row in enumerate(rows):
        minimal_table_row(s, 2.08 + i * 0.80, widths, row, h=0.80,
                          sizes=(10.2, 10.0, 9.8, 9.8))
    s.text(0.78, 6.38, 11.45, 0.34,
           "No worker is running. No selection, training, validation decoding or test access has occurred.",
           size=12.5, bold=True, font="sans", text_color=INK, align="center")
    minimal_footer(s); slides.append(s)

    # 10 — graphical Gantt chart
    s = SlideScene(10, "Time plan", "Project Gantt chart",
                   ("Time Plan",),
                   "Detailed dependencies and acceptance evidence: docs/gantt-plan.md")
    minimal_header(s, s.title)
    chart_x0, chart_x1 = 3.38, 12.24
    chart_y0, row_h = 2.08, 0.54
    start_date, end_date = date(2026, 8, 9), date(2026, 11, 22)
    total_days = (end_date - start_date).days

    def gx(day: date) -> float:
        return chart_x0 + ((day - start_date).days / total_days) * (chart_x1 - chart_x0)

    # Month bands and exact date axis.
    months = [
        (date(2026, 8, 9), date(2026, 9, 1), "AUG"),
        (date(2026, 9, 1), date(2026, 10, 1), "SEP"),
        (date(2026, 10, 1), date(2026, 11, 1), "OCT"),
        (date(2026, 11, 1), end_date, "NOV"),
    ]
    s.text(0.78, 1.45, 2.36, 0.26, "Workstream", size=11.0, bold=True,
           font="sans", text_color=INK)
    for month_start, month_end, label_text in months:
        x0, x1 = gx(month_start), gx(month_end)
        s.rect(x0, 1.42, x1 - x0, 0.42, fill="F3F3F3", stroke=LINE,
               stroke_width=0.5, radius=0)
        s.text(x0, 1.54, x1 - x0, 0.16, label_text, size=9.3, bold=True,
               font="sans", text_color=INK, align="center")
        s.line(x0, 1.84, 0, 4.28, stroke=LINE, stroke_width=0.6)
    s.line(chart_x1, 1.84, 0, 4.28, stroke=LINE, stroke_width=0.6)

    gantt_rows = [
        ("BLER characterization + G-8", date(2026, 8, 9), date(2026, 8, 23), "IN PROGRESS", "555555", WHITE),
        ("Learned model training + tuning", date(2026, 8, 24), date(2026, 9, 21), "PLANNED", "D9D9D9", INK),
        ("Digital feature system + validation", date(2026, 9, 21), date(2026, 10, 5), "PLANNED", "D9D9D9", INK),
        ("Freeze + single test campaign", date(2026, 10, 5), date(2026, 10, 19), "PLANNED", "D9D9D9", INK),
        ("Demo + optional replay", date(2026, 10, 19), date(2026, 11, 2), "PLANNED", "D9D9D9", INK),
        ("Report, audit + final review", date(2026, 11, 2), date(2026, 11, 22), "PLANNED", "D9D9D9", INK),
    ]
    for i, (label_text, bar_start, bar_end, state, fill, text_ink) in enumerate(gantt_rows):
        y = chart_y0 + i * row_h
        s.text(0.78, y + 0.08, 2.42, 0.24, label_text, size=10.7,
               font="sans", text_color=INK, valign="mid")
        s.line(0.78, y + row_h - 0.03, 11.46, 0, stroke=LINE, stroke_width=0.35)
        x0, x1 = gx(bar_start), gx(bar_end)
        s.rect(x0, y + 0.08, max(0.18, x1 - x0), 0.32, fill=fill,
               stroke="888888", stroke_width=0.5, radius=0)
        if x1 - x0 > 1.05:
            s.text(x0 + 0.05, y + 0.16, x1 - x0 - 0.10, 0.13, state,
                   size=7.1, bold=True, font="sans", text_color=text_ink,
                   align="center", valign="mid")

    # Fixed review/report dates are markers, not movable task bars.
    markers = [
        (date(2026, 8, 18), "R1"),
        (date(2026, 9, 29), "R2"),
        (date(2026, 11, 17), "FINAL"),
        (date(2026, 11, 20), "REPORT"),
    ]
    for marker_date, marker_label in markers:
        x = gx(marker_date)
        s.line(x, 1.84, 0, 4.10, stroke="777777", stroke_width=0.8)
        s.circle(x - 0.055, 1.89, 0.11, 0.11, fill=INK, stroke=None)

    # W16 is deliberately reserved inside the reporting bar.
    contingency_x0, contingency_x1 = gx(date(2026, 11, 9)), gx(date(2026, 11, 16))
    contingency_y = chart_y0 + 5 * row_h
    s.rect(contingency_x0, contingency_y + 0.04,
           contingency_x1 - contingency_x0, 0.40, fill=None,
           stroke=INK, stroke_width=1.1, radius=0)
    s.text(contingency_x0, contingency_y + 0.47,
           contingency_x1 - contingency_x0, 0.16, "W16 contingency",
           size=6.7, font="sans", text_color=MUTED, align="center")

    s.text(0.78, 6.20, 11.45, 0.24,
           "Review dates: 18–22 Aug · 29 Sep–3 Oct · 17–21 Nov · report due 20 Nov",
           size=11.4, font="sans", text_color=INK, align="center")
    s.text(0.78, 6.54, 11.45, 0.22,
           "BLER characterization is the only workstream in progress. All later bars are planned; W16 is report contingency.",
           size=10.4, font="sans", text_color=MUTED, align="center")
    minimal_footer(s); slides.append(s)

    # 11 — application, scope and risks expected in a preliminary review
    s = SlideScene(11, "Application and scope", "Applications, project scope and risks",
                   ("Motivation", "Subject Knowledge", "Objectives"),
                   "Circular clause 5 · PR-3 standards register · PR-9 deployment dossier")
    minimal_header(s, s.title)
    columns = [
        (0.78, "Edge AI application",
         "Remote cameras and edge sensors send task-relevant information to a receiver-side AI model when bandwidth, latency or link reliability is limited. This represents split edge/cloud inference."),
        (4.42, "AI/ML contribution",
         "Channel-robust neural representation learning; a residual encoder and dual-head decoder; CE + λMSE multi-task training; and a digital feature control that isolates the value of joint coding."),
        (8.06, "Main risks and controls",
         "BLER work is incomplete: resume from the authenticated checkpoint.\n\nA hypothesis may not be supported: report the result without retuning.\n\nHardware adds confounders: keep it as gated stretch work."),
    ]
    for x, heading, body in columns:
        s.rect(x, 1.52, 3.34, 3.62, fill=WHITE, stroke=LINE, stroke_width=0.8, radius=0)
        s.text(x + 0.20, 1.78, 2.94, 0.56, heading, size=15.2,
               bold=True, font="sans", text_color=INK)
        s.text(x + 0.20, 2.48, 2.94, 2.30, body, size=12.5,
               font="sans", text_color=MUTED)
    s.text(0.80, 5.52, 2.54, 0.30, "Technical claim boundary", size=14.2,
           bold=True, font="sans", text_color=INK)
    s.text(3.38, 5.47, 8.76, 0.58,
           "Tier 1 is a simulation-based AI experiment over normalized AWGN. OpenJPEG 2.5.4 and TS 38.212-derived LDPC are controls; no complete 5G NR radio or required hardware is claimed.",
           size=12.7, font="sans", text_color=MUTED)
    s.rect(0.78, 6.24, 11.40, 0.46, fill="F5F5F5", stroke=None, radius=0)
    s.text(1.02, 6.36, 10.92, 0.20,
           "Outstanding human item: the guide's dated acknowledgement of the no-required-hardware path remains pending.",
           size=11.8, bold=True, font="sans", text_color=INK, align="center")
    minimal_footer(s); slides.append(s)

    # 12 — close with the project summary, next steps and rubric map
    s = SlideScene(12, "Review summary", "Summary and next steps",
                   ("Motivation", "Objectives", "Hypothesis", "Problem Survey", "Subject Knowledge", "Time Plan"),
                   "Full package: deliverables/review-1/first-review-package.md")
    minimal_header(s, s.title)
    s.text(0.78, 1.42, 6.74, 0.30, "Project summary", size=15.5,
           bold=True, font="sans", text_color=INK)
    summary_rows = [
        ("Problem", "Edge AI inference must operate over a limited noisy link."),
        ("AI/ML method", "Train a residual neural encoder and dual-head decoder through a differentiable channel."),
        ("Evaluation", "Measure Imagenette-160 classification at matched channel use against two controlled baselines."),
        ("Next work", "Finish characterization, train and freeze the models, then run the paired test campaign once."),
    ]
    for i, (label_text, body) in enumerate(summary_rows):
        y = 1.92 + i * 0.82
        s.text(0.80, y, 1.18, 0.28, label_text, size=13.2, bold=True,
               font="sans", text_color=INK)
        s.text(2.04, y - 0.02, 5.42, 0.56, body, size=12.7,
               font="sans", text_color=MUTED)
        s.line(0.80, y + 0.60, 6.66, 0, stroke=LINE, stroke_width=0.4)
    s.line(7.78, 1.42, 0, 4.56, stroke=LINE, stroke_width=0.7)
    s.text(8.08, 1.42, 4.02, 0.30, "Six-criterion rubric map", size=15.5,
           bold=True, font="sans", text_color=INK)
    rubric = [
        ("Motivation", "2–4, 6, 11"),
        ("Objectives", "6–7, 11–12"),
        ("Hypothesis", "8"),
        ("Problem Survey", "5"),
        ("Subject Knowledge", "3–5, 7–9, 11"),
        ("Time Plan", "10"),
    ]
    for i, (criterion, slide_refs) in enumerate(rubric):
        y = 1.94 + i * 0.58
        s.text(8.08, y, 2.22, 0.28, criterion, size=12.0, bold=True,
               font="sans", text_color=INK)
        s.text(10.36, y, 1.72, 0.28, slide_refs, size=12.0,
               font="sans", text_color=MUTED, align="right")
        s.line(8.08, y + 0.36, 4.00, 0, stroke=LINE, stroke_width=0.4)
    s.rect(0.78, 5.52, 11.30, 1.02, fill="F5F5F5", stroke=None, radius=0)
    s.text(1.02, 5.72, 2.02, 0.26, "First Review position", size=13.0,
           bold=True, font="sans", text_color=INK)
    s.text(3.04, 5.67, 8.76, 0.58,
           "The AI/ML method, controlled evaluation protocol, feasibility evidence and completion plan are defined. Final neural training and learned-versus-classical results remain future work.",
           size=13.1, font="sans", text_color=INK)
    minimal_footer(s); slides.append(s)

    return slides


def ppt_font(el: Element) -> str:
    return {"serif": TITLE_FONT, "sans": BODY_FONT, "math": MATH_FONT, "mono": MONO_FONT}.get(el.font, BODY_FONT)


def add_to_ppt(slide, el: Element) -> None:
    if el.kind == "text":
        box = slide.shapes.add_textbox(Inches(el.x), Inches(el.y), Inches(el.w), Inches(el.h))
        tf = box.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.margin_left = Inches(el.margin)
        tf.margin_right = Inches(el.margin)
        tf.margin_top = Inches(el.margin)
        tf.margin_bottom = Inches(el.margin)
        tf.vertical_anchor = {"top": MSO_ANCHOR.TOP, "mid": MSO_ANCHOR.MIDDLE, "bottom": MSO_ANCHOR.BOTTOM}.get(el.valign, MSO_ANCHOR.TOP)
        p = tf.paragraphs[0]
        p.text = el.text
        p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}.get(el.align, PP_ALIGN.LEFT)
        p.space_before = Pt(0)
        p.space_after = Pt(0)
        p.line_spacing = 1.0
        for run in p.runs:
            run.font.name = ppt_font(el)
            run.font.size = Pt(el.size)
            run.font.bold = el.bold
            run.font.italic = el.italic
            run.font.color.rgb = rgb(el.text_color)
        box.rotation = el.rotation
    elif el.kind in {"rect", "circle"}:
        shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if el.kind == "rect" and el.radius else MSO_SHAPE.RECTANGLE
        if el.kind == "circle":
            shape_type = MSO_SHAPE.OVAL
        shape = slide.shapes.add_shape(shape_type, Inches(el.x), Inches(el.y), Inches(el.w), Inches(el.h))
        if el.fill:
            shape.fill.solid(); shape.fill.fore_color.rgb = rgb(el.fill)
        else:
            shape.fill.background()
        if el.stroke:
            shape.line.color.rgb = rgb(el.stroke); shape.line.width = Pt(el.stroke_width)
        else:
            shape.line.fill.background()
    elif el.kind in {"line", "arrow"}:
        line = slide.shapes.add_connector(1, Inches(el.x), Inches(el.y), Inches(el.x + el.w), Inches(el.y + el.h))
        line.line.color.rgb = rgb(el.stroke or INK)
        line.line.width = Pt(el.stroke_width)
        if el.kind == "arrow":
            # python-pptx does not expose arrowheads consistently; use a small native triangle.
            angle = math.atan2(el.h, el.w) if el.w or el.h else 0
            tip_x, tip_y = el.x + el.w, el.y + el.h
            tri = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE,
                                         Inches(tip_x - 0.07), Inches(tip_y - 0.055),
                                         Inches(0.14), Inches(0.11))
            tri.rotation = math.degrees(angle) + 90
            tri.fill.solid(); tri.fill.fore_color.rgb = rgb(el.stroke or INK)
            tri.line.fill.background()


def build_pptx(scenes: list[SlideScene]) -> None:
    prs = Presentation()
    prs.slide_width = Inches(SW)
    prs.slide_height = Inches(SH)
    prs.core_properties.title = "Task-Oriented Deep Joint Source–Channel Coding — First Review"
    prs.core_properties.subject = "Capstone First Review, 18–22 August 2026"
    prs.core_properties.author = "Capstone project team"
    prs.core_properties.comments = "Generated from authenticated repository evidence; editable native shapes and text."
    blank = prs.slide_layouts[6]
    for scene in scenes:
        slide = prs.slides.add_slide(blank)
        bg = slide.background.fill
        bg.solid(); bg.fore_color.rgb = rgb(IVORY)
        for el in scene.elements:
            add_to_ppt(slide, el)
    prs.save(PPTX_PATH)


def pil_font(el: Element) -> ImageFont.FreeTypeFont:
    if el.font in {"serif", "math"}:
        key = "serif_bold" if el.bold else "serif"
    elif el.font == "mono":
        key = "mono"
    else:
        key = "sans_bold" if el.bold else "sans"
    return ImageFont.truetype(FONT_FILES[key], max(8, round(el.size * PX_H / 540)))


def wrap_for_draw(draw: ImageDraw.ImageDraw, text_value: str, font, max_width: int) -> str:
    wrapped: list[str] = []
    for raw in text_value.split("\n"):
        if not raw:
            wrapped.append("")
            continue
        words = raw.split()
        line = ""
        for word in words:
            candidate = word if not line else f"{line} {word}"
            if draw.textbbox((0, 0), candidate, font=font)[2] <= max_width:
                line = candidate
            else:
                if line:
                    wrapped.append(line)
                line = word
        if line:
            wrapped.append(line)
    return "\n".join(wrapped)


def render_element(draw: ImageDraw.ImageDraw, el: Element) -> None:
    x, y, w, h = px(el.x, "x"), px(el.y, "y"), px(el.w, "x"), px(el.h, "y")
    if el.kind == "rect":
        r = max(0, px(el.radius, "y"))
        draw.rounded_rectangle((x, y, x + w, y + h), radius=r,
                               fill=color(el.fill) if el.fill else None,
                               outline=color(el.stroke) if el.stroke else None,
                               width=max(1, round(el.stroke_width * 1.5)))
    elif el.kind == "circle":
        draw.ellipse((x, y, x + w, y + h), fill=color(el.fill) if el.fill else None,
                     outline=color(el.stroke) if el.stroke else None,
                     width=max(1, round(el.stroke_width * 1.5)))
    elif el.kind in {"line", "arrow"}:
        x2, y2 = px(el.x + el.w, "x"), px(el.y + el.h, "y")
        draw.line((x, y, x2, y2), fill=color(el.stroke or INK), width=max(1, round(el.stroke_width * 2)))
        if el.kind == "arrow":
            angle = math.atan2(y2 - y, x2 - x)
            length = 11
            pts = [(x2, y2),
                   (x2 - length * math.cos(angle - 0.48), y2 - length * math.sin(angle - 0.48)),
                   (x2 - length * math.cos(angle + 0.48), y2 - length * math.sin(angle + 0.48))]
            draw.polygon(pts, fill=color(el.stroke or INK))
    elif el.kind == "text":
        font = pil_font(el)
        margin_x, margin_y = px(el.margin, "x"), px(el.margin, "y")
        max_width = max(4, w - 2 * margin_x)
        rendered = wrap_for_draw(draw, el.text, font, max_width)
        bbox = draw.multiline_textbbox((0, 0), rendered, font=font, spacing=2, align=el.align)
        tw, th = bbox[2] - bbox[0], bbox[3] - bbox[1]
        tx = x + margin_x
        if el.align == "center": tx = x + (w - tw) / 2
        elif el.align == "right": tx = x + w - margin_x - tw
        ty = y + margin_y
        if el.valign == "mid": ty = y + (h - th) / 2
        elif el.valign == "bottom": ty = y + h - margin_y - th
        draw.multiline_text((tx, ty), rendered, font=font, fill=color(el.text_color),
                            spacing=2, align=el.align)


def render_previews(scenes: list[SlideScene]) -> list[Path]:
    if PREVIEWS.exists():
        shutil.rmtree(PREVIEWS)
    PREVIEWS.mkdir(parents=True)
    paths: list[Path] = []
    images: list[Image.Image] = []
    for scene in scenes:
        img = Image.new("RGB", (PX_W, PX_H), color(IVORY))
        draw = ImageDraw.Draw(img)
        for el in scene.elements:
            render_element(draw, el)
        path = PREVIEWS / f"slide-{scene.number:02d}.png"
        img.save(path, quality=95)
        paths.append(path)
        images.append(img)
    images[0].save(PDF_PATH, "PDF", resolution=144.0, save_all=True,
                   append_images=images[1:])
    # 4×3 contact sheet for fast review.
    thumb_w, thumb_h = 533, 300
    sheet = Image.new("RGB", (thumb_w * 3, thumb_h * 4), color("D6D3CC"))
    for i, img in enumerate(images):
        thumb = img.resize((thumb_w - 8, thumb_h - 8), Image.Resampling.LANCZOS)
        sheet.paste(thumb, ((i % 3) * thumb_w + 4, (i // 3) * thumb_h + 4))
    sheet.save(CONTACT_PATH, quality=95)
    return paths


def main() -> None:
    OUT.mkdir(parents=True, exist_ok=True)
    scenes = build_minimal_scenes()
    assert len(scenes) == 12
    build_pptx(scenes)
    paths = render_previews(scenes)
    print(f"wrote {PPTX_PATH}")
    print(f"wrote {PDF_PATH}")
    print(f"wrote {CONTACT_PATH}")
    print(f"wrote {len(paths)} slide previews under {PREVIEWS}")


if __name__ == "__main__":
    main()
